"""
PPTX Generator 
日本企業向けハイデザインプレゼンテーション生成システム

【使い方】3ステップで資料作成
  STEP 1: SLIDES_TO_USE でスライド構成を決める
    使いたいテンプレート番号を配列で指定。順序が反映され、同じ番号の複数回使用OK。
    
    例: SLIDES_TO_USE = [1, 2, 3, 4, 5, 1]
        → 表紙(1) → 目次(2) → 棒グラフ(3) → 円グラフ(4) → テーブル(5) → 表紙(1)

  STEP 2: SLIDE_CONTENT でテキスト内容を編集
    キー = 出力順（1始まり）。SLIDES_TO_USE の順序に対応。
    
    例: SLIDE_CONTENT = {
            1: {"title": "表紙タイトル", "subtitle": "...", "credit": "..."},
            2: {"title": "目次", "items": ["項目1", "項目2", ...]},
            3: {"title": "売上推移", "categories": [...], "series": [...], ...},
            ...
        }
  STEP 3: 実行
    python pptx_generator.py
    → outputs/pptx_generator.pptx が生成される

================================================================================
【テンプレート一覧】全12種

  ■ 可変コンテンツ対応（items/categories/rows等の配列サイズに応じて自動調整）
    1: 表紙                    → title, subtitle, credit
    2: 目次（リスト）           → items[] ★可変: 1〜10項目程度OK
    3: 左グラフ・右テキスト     → categories[], series[], items[] ★可変
    4: 左テキスト・右グラフ     → categories[], series[], items[] ★可変
    5: 左テーブル・右テキスト   → columns[], rows[][] ★可変: 行列数自由
    6: 左テキスト・右グラフ     → categories[], series[], items[] ★可変
    7: カードグリッド（4枚）    → cards[] ★可変: 1〜4枚（4枚固定レイアウト）
    8: 左テキスト・右画像       → items[] ★可変, image_url

  ■ 固定レイアウト（項目数を変えるとレイアウト崩れの可能性あり）
    9:  3ポイント（丸囲み）     → items[3] ⚠固定: 3項目推奨
    10: ポイントリスト（コンパクト） → items[:4] ⚠上限: 最大4項目
    11: ポイントリスト（大）    → items[:3] ⚠上限: 最大3項目
    12: タイムライン           → items[:4] ⚠固定: 円+バー配置のため変更注意

================================================================================
【データ構造リファレンス】

  ● 箇条書き（items）
    "items": ["項目1", "項目2", "項目3"]
    → 配列の長さに応じて行数が変わる

  ● グラフ（chart）
    "categories": ["Q1", "Q2", "Q3", "Q4"],
    "series": [
        {"name": "2024年", "values": [100, 120, 140, 160]},
        {"name": "2025年", "values": [110, 130, 150, 170]}
    ],
    "chart_type": "COLUMN_CLUSTERED"  # or "PIE", "LINE"
    → categoriesの数、seriesの数は自由に増減可能

  ● テーブル（table）
    "columns": ["項目", "2023年", "2024年", "2025年"],
    "rows": [
        ["売上", "100億", "120億", "140億"],
        ["利益", "10億", "15億", "20億"]
    ]
    → 列数・行数は自由に増減可能

  ● カード（cards）
    "cards": [
        {"icon": "🚀", "title": "タイトル1", "desc": "説明1"},
        {"icon": "💡", "title": "タイトル2", "desc": "説明2"}
    ]

================================================================================
【ファイル構成】編集箇所の目安

-セクション1: テーマ定義 (Themes)                    ← デザイン（色・フォント）
-セクション2: テンプレート定義 (Templates)           ← デザイン（レイアウト）  
-セクション3: データ定義 (Configuration)             ← ★ここを編集 ★         
 -SLIDE_TEMPLATES: テンプレート一覧                                      
 -SLIDES_TO_USE: 使用するスライド番号             ← ★ここを編集 ★         
 -SLIDE_CONTENT: 各スライドのテキスト             ← ★ここを編集 ★         
-セクション4: エンジン & 実行                        ← 触らなくてOK            


"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls


# ============================================================
# 1. テーマ定義 (Themes) ★デザイン: 色・フォント★
# ============================================================
# ※通常は編集不要。ブランドカラーを変更したい場合のみ編集。
class ModernJapanTheme:
    """
    シンプル＆洗練された青×白カラースキーム
    参考: スライドデザイン帳スタイル
    """
    # Primary: コーポレートブルー（濃い青）
    PRIMARY = RGBColor(0, 82, 155)  # #00529B
    
    # Secondary: アクセントブルー（明るい青）
    SECONDARY = RGBColor(0, 114, 198)  # #0072C6
    
    # Accent: ポイントカラー（さらに明るい青）
    ACCENT = RGBColor(65, 143, 222)  # #418FDE
    
    # Chart Colors (Blue Gradation Palette)
    CHART_COLORS = [
        RGBColor(0, 82, 155),    # Primary Blue (濃)
        RGBColor(0, 114, 198),   # Secondary Blue
        RGBColor(65, 143, 222),  # Accent Blue
        RGBColor(130, 180, 230), # Light Blue
        RGBColor(180, 210, 240), # Lighter Blue
        RGBColor(220, 235, 250)  # Very Light Blue
    ]
    
    # Base Colors
    BACKGROUND = RGBColor(255, 255, 255)  # 白
    TEXT = RGBColor(51, 51, 51)           # 濃いグレー #333333
    TEXT_LIGHT = RGBColor(102, 102, 102)  # ミディアムグレー #666666
    WHITE = RGBColor(255, 255, 255)
    
    # UI Elements
    DIVIDER_LINE = RGBColor(0, 82, 155)   # 青のライン
    BOX_BORDER = RGBColor(200, 220, 240)  # 薄い青のボーダー
    LIGHT_BG = RGBColor(240, 247, 255)    # 薄い青の背景
    
    # Number/Point Colors
    POINT_NUMBER = RGBColor(0, 82, 155)   # ポイント番号の青
    
    # Fonts
    FONT_JP = "メイリオ"
    FONT_EN = "Segoe UI"


# 使用するテーマを選択
CURRENT_THEME = ModernJapanTheme()


# ============================================================
# Helper: OXML Bullet Injection
# ============================================================
def _enable_bullet(paragraph):
    """段落に対してPowerPoint本来の箇条書き属性を付与"""
    pPr = paragraph._p.get_or_add_pPr()
    buChar = parse_xml('<a:buChar %s char="•"/>' % nsdecls('a'))
    if pPr.find(f"{{{nsdecls('a')}}}buChar") is None:
        pPr.insert(0, buChar)
    pPr.set('marL', '457200')
    pPr.set('indent', '-228600')


# ============================================================
# 2. テンプレート定義 (Templates) ★デザイン: レイアウト★
# ============================================================
# ※通常は編集不要。スライドのレイアウトを変更したい場合のみ編集。
def create_cover(slide, data, theme):
    """表紙 - 青×白のシンプルデザイン"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme.BACKGROUND

    # 左側の青いバー（統一: 0.08インチ）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.PRIMARY
    bar.line.fill.background()
    
    # 装飾: 右下の三角形
    triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(8.5), Inches(4.0), Inches(1.5), Inches(1.5))
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = theme.PRIMARY
    triangle.line.fill.background()
    triangle.rotation = 90

    title = data.get('title', '')
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = theme.FONT_JP
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme.PRIMARY

    if 'subtitle' in data:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.0), Inches(0.5))
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = theme.FONT_EN
        p.font.size = Pt(14)
        p.font.color.rgb = theme.SECONDARY
        p.font.bold = True

    if 'credit' in data:
        credit_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(8.0), Inches(0.5))
        p = credit_box.text_frame.paragraphs[0]
        p.text = data['credit']
        p.font.name = theme.FONT_JP
        p.font.size = Pt(12)
        p.font.color.rgb = theme.TEXT_LIGHT


def create_list(slide, data, theme):
    """リスト - シンプル箇条書き"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.name = theme.FONT_JP
        p.font.size = Pt(16)
        p.font.color.rgb = theme.TEXT
        p.space_after = Pt(14)
        _enable_bullet(p)


def create_chart_left_text_right(slide, data, theme):
    """左：グラフ / 右：テキスト"""
    _setup_header(slide, data, theme)
    _create_chart_at(slide, data, theme, Inches(0.5), Inches(1.4), Inches(5.5), Inches(3.8))
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(6.2), Inches(1.4), Inches(3.3))


def create_table_left_text_right(slide, data, theme):
    """左：テーブル / 右：テキスト"""
    _setup_header(slide, data, theme)
    _create_table_at(slide, data, theme, Inches(0.5), Inches(1.4), Inches(5.5))
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(6.2), Inches(1.4), Inches(3.3))


def create_text_left_chart_right(slide, data, theme):
    """左：テキスト / 右：グラフ"""
    _setup_header(slide, data, theme)
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(0.5), Inches(1.4), Inches(3.5))
    _create_chart_at(slide, data, theme, Inches(4.2), Inches(1.4), Inches(5.3), Inches(3.8))


def create_text_left_image_right(slide, data, theme):
    """左：テキスト / 右：画像プレースホルダー"""
    _setup_header(slide, data, theme)
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(0.5), Inches(1.4), Inches(4.5))
    
    # 画像プレースホルダーを表示
    _add_image_placeholder(slide, theme)


def _add_image_placeholder(slide, theme):
    """画像プレースホルダーを追加（フォールバック用）"""
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.3), Inches(3.5))
    frame.fill.background()
    frame.line.color.rgb = theme.BOX_BORDER
    frame.line.dash_style = 1
    
    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.9), Inches(4.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "Image Area"
    p.font.size = Pt(12)
    p.font.color.rgb = theme.TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER


def create_card_grid(slide, data, theme):
    """カードグリッド - シンプル4カードスタイル"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    start_y = Inches(1.4)
    
    for i, item in enumerate(items[:4]):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(4.6)
        y = start_y + row * Inches(2.0)
        w = Inches(4.2)
        h = Inches(1.8)
        
        # カード背景
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = theme.WHITE
        card.line.color.rgb = theme.BOX_BORDER
        
        # 上部の青いライン
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.04))
        bar.fill.solid()
        bar.fill.fore_color.rgb = theme.PRIMARY
        bar.line.fill.background()

        # テキスト
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), h - Inches(0.4))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = item.replace('\n', '\n')
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.color.rgb = theme.TEXT
        p.alignment = PP_ALIGN.LEFT


def create_three_points_circle(slide, data, theme):
    """3ポイント - 丸囲み番号スタイル"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    
    start_y = Inches(1.5)
    col_width = Inches(2.8)
    gap = Inches(0.4)
    start_x = Inches(0.7)
    
    for i, item in enumerate(items[:3]):
        x = start_x + i * (col_width + gap)
        
        # 丸い枠（円）
        circle_x = x + Inches(0.8)
        circle_y = start_y
        circle_size = Inches(1.4)
        
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, circle_size, circle_size)
        circle.fill.background()
        circle.line.color.rgb = theme.PRIMARY
        circle.line.width = Pt(2)
        
        # 番号（円の中央 - 手動で位置調整）
        # 円の中心: circle_y + circle_size/2 = start_y + 0.7インチ
        # テキスト高さ約0.5インチとして、中央に配置
        text_height = Inches(0.55)
        text_y = circle_y + (circle_size - text_height) / 2 + Inches(0.02)  # 微調整
        
        num_box = slide.shapes.add_textbox(circle_x, text_y, circle_size, text_height)
        tf = num_box.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        
        p = tf.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        p.alignment = PP_ALIGN.CENTER  # 水平中央揃え
        
        # 見出しと説明文を分離
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        # 見出し（円の下）
        head_box = slide.shapes.add_textbox(x, start_y + Inches(1.55), col_width, Inches(0.4))
        head_box.text_frame.word_wrap = True
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # 説明文
        if desc:
            desc_box = slide.shapes.add_textbox(x, start_y + Inches(1.95), col_width, Inches(1.5))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(10)
            p.font.color.rgb = theme.TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER


def create_point_list_compact(slide, data, theme):
    """ポイントリスト - コンパクト（番号 + 見出し + 説明）"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    
    start_y = Inches(1.4)
    item_height = Inches(0.85)
    
    for i, item in enumerate(items[:4]):
        y_pos = start_y + i * item_height
        
        # 番号
        num_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(0.5), Inches(0.5))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        
        # 見出しと説明文を分離
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        # 見出し
        head_box = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(8.0), Inches(0.35))
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        
        # 説明文
        if desc:
            desc_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.35), Inches(8.0), Inches(0.45))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(11)
            p.font.color.rgb = theme.TEXT_LIGHT


def create_point_list_large(slide, data, theme):
    """ポイントリスト - 大きめ（左サイドバー + 番号 + テキスト）"""
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme.BACKGROUND
    
    # 左側の青いバー（統一: 0.08インチ）
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = theme.PRIMARY
    left_bar.line.fill.background()
    
    items = data.get('items', [])
    start_y = Inches(0.5)
    item_height = Inches(1.7)
    
    for i, item in enumerate(items[:3]):
        y_pos = start_y + i * item_height
        
        # 番号（大きく）
        num_box = slide.shapes.add_textbox(Inches(0.4), y_pos, Inches(1.0), Inches(1.0))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        
        # 見出しと説明文を分離
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        # 見出し
        head_box = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.1), Inches(7.8), Inches(0.4))
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        
        # 説明文
        if desc:
            desc_box = slide.shapes.add_textbox(Inches(1.6), y_pos + Inches(0.55), Inches(7.8), Inches(1.0))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(12)
            p.font.color.rgb = theme.TEXT_LIGHT
        
        # 区切り線（最後以外）
        if i < len(items[:3]) - 1:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), y_pos + item_height - Inches(0.1), Inches(9.2), Inches(0.01))
            line.fill.solid()
            line.fill.fore_color.rgb = theme.DIVIDER_LINE
            line.line.fill.background()


# --- Helper Functions ---

def _create_text_list_bullet(slide, items, theme, x, y, w):
    """バレット付きテキストリスト - シンプル箇条書き"""
    tb = slide.shapes.add_textbox(x, y, w, Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.color.rgb = theme.TEXT
        p.space_after = Pt(10)
        _enable_bullet(p)


def _create_chart_at(slide, data, theme, x, y, w, h):
    """グラフ作成（色設定・円グラフ対応）"""
    
    chart_type_str = data.get('chart_type', 'COLUMN_CLUSTERED')
    
    # 円グラフ（Pie Chart）の場合のデータ構造は異なる
    if chart_type_str == 'PIE':
        chart_data = ChartData()
        chart_data.categories = data.get('categories', [])
        # 円グラフは最初のシリーズのみ使用
        series_list = data.get('series', [])
        if series_list:
            chart_data.add_series(series_list[0]['name'], series_list[0]['values'])
        chart_type = XL_CHART_TYPE.PIE
    else:
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', [])
        series_list = data.get('series', [])
        for s in series_list:
            chart_data.add_series(s['name'], s['values'])
        chart_type = getattr(XL_CHART_TYPE, chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # チャート追加
    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart
    
    # スタイル・凡例設定
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = theme.FONT_JP
    chart.chart_title.has_text_frame = False
    
    # データラベル（円グラフ用）
    if chart_type == XL_CHART_TYPE.PIE:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(9)
        data_labels.font.color.rgb = theme.TEXT
    
    # シリーズの色設定
    # python-pptxのchart.seriesのAPIは限定的だが、可能な範囲で適用
    try:
        if chart_type == XL_CHART_TYPE.PIE:
            # 円グラフは「ポイント」ごとに色を変える
            plot = chart.plots[0]
            series = plot.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                # 用意したパレットを循環
                color = theme.CHART_COLORS[i % len(theme.CHART_COLORS)]
                fill.fore_color.rgb = color
        else:
            # 棒・折れ線は「シリーズ」ごとに色を変える
            for i, series in enumerate(chart.series):
                fill = series.format.fill
                fill.solid()
                color = theme.CHART_COLORS[i % len(theme.CHART_COLORS)]
                fill.fore_color.rgb = color
                
                # 折れ線グラフの場合、線自体の色も変える必要があるが
                # line.colorプロパティへのアクセスはバージョンにより異なる場合あり
                if hasattr(series.format.line, 'color'):
                     series.format.line.color.rgb = color

    except Exception as e:
        # 色設定に失敗してもグラフ自体は生成されるようにする
        print(f"Note: Color formatting skipped for some elements ({e})")


def _create_table_at(slide, data, theme, x, y, w):
    """テーブル - 青ヘッダースタイル"""
    columns = data.get('columns', [])
    rows = data.get('rows', [])
    rows_count = len(rows) + 1
    cols_count = len(columns)
    
    row_height = Inches(0.4)
    table_height = row_height * rows_count
    
    shape = slide.shapes.add_table(rows_count, cols_count, x, y, w, table_height)
    table = shape.table
    
    # ヘッダー行（青背景・白文字）
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.PRIMARY  # 青背景
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE  # 白文字
        p.alignment = PP_ALIGN.CENTER

    # データ行
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            # 交互に薄い青背景
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.LIGHT_BG
            else:
                cell.fill.background()
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = theme.TEXT
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

def _setup_header(slide, data, theme):
    """ヘッダー - 青ライン＋タイトル"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme.BACKGROUND

    # 左側の青いバー（細め）
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = theme.PRIMARY
    left_bar.line.fill.background()

    # タイトル下の青いライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.PRIMARY
    line.line.fill.background()

    if 'subtitle' in data:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = theme.FONT_EN
        p.font.size = Pt(10)
        p.font.color.rgb = theme.SECONDARY
        p.font.bold = True

    title = data.get('title', '')
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(8), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = theme.FONT_JP
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = theme.PRIMARY  # タイトルを青に


# ------------------------------------------------------------
# タイムラインスライド
# ------------------------------------------------------------
def create_timeline(slide, data, theme):
    """
    タイムラインスライド
    - 横長のカラーバー（最大4セグメント）
    - 各ポイントに年/ラベル + 説明
    - 上下交互に配置（重ならないように調整）
    """
    # 背景色
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = theme.BACKGROUND
    
    # 左サイドバー
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(5.625))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = theme.PRIMARY
    sidebar.line.fill.background()
    
    # ヘッダー（シンプル版 - タイトルのみ）
    if 'subtitle' in data and data['subtitle']:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = theme.FONT_EN
        p.font.size = Pt(10)
        p.font.color.rgb = theme.SECONDARY
        p.font.bold = True
    
    title = data.get('title', '')
    if title:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(8), Inches(0.5))
        p = title_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = theme.FONT_JP
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
    
    # タイムラインデータ
    timeline_items = data.get('timeline_items', [
        {"year": "2021", "title": "Phase 1", "description": "説明文をここに入力"},
        {"year": "2022", "title": "Phase 2", "description": "説明文をここに入力"},
        {"year": "2023", "title": "Phase 3", "description": "説明文をここに入力"},
        {"year": "2024", "title": "Phase 4", "description": "説明文をここに入力"},
    ])
    
    # タイムラインの色（青系グラデーション）
    timeline_colors = [
        RGBColor(0, 82, 155),    # PRIMARY - 濃い青
        RGBColor(0, 114, 198),   # SECONDARY - 明るい青
        RGBColor(70, 150, 200),  # ライトブルー
        RGBColor(100, 180, 220), # より明るいブルー
    ]
    
    num_items = min(len(timeline_items), 4)  # 最大4項目
    
    # タイムラインバーの位置（中央に配置）
    bar_y = Inches(2.7)
    bar_height = Inches(0.45)
    bar_start_x = Inches(0.5)
    bar_total_width = Inches(9.0)
    segment_width = bar_total_width / num_items
    
    # 円のサイズ
    circle_size = Inches(0.5)
    
    # 各セグメントを描画
    for i in range(num_items):
        item = timeline_items[i]
        seg_x = bar_start_x + segment_width * i
        
        # セグメントの矩形
        seg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            seg_x, bar_y, 
            segment_width, bar_height
        )
        seg.fill.solid()
        seg.fill.fore_color.rgb = timeline_colors[i % len(timeline_colors)]
        seg.line.fill.background()
        
        # 年ラベル（バー内）
        year_box = slide.shapes.add_textbox(seg_x, bar_y, segment_width, bar_height)
        tf = year_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = item.get('year', '')
        p.font.name = theme.FONT_EN
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE
        
        # 上下交互に配置
        is_top = (i % 2 == 0)  # 偶数は上、奇数は下
        
        # 縦線の位置（セグメントの中央）
        line_x = seg_x + segment_width / 2
        
        if is_top:
            # 上側に配置
            circle_y = Inches(1.65)  # 円の位置（バーから離す）
            line_start_y = circle_y + circle_size  # 円の下から
            line_end_y = bar_y  # バーの上まで
            title_y = Inches(1.05)  # タイトル位置
            desc_y = Inches(1.30)   # 説明文位置
        else:
            # 下側に配置
            circle_y = Inches(3.7)  # 円の位置（バーから離す）
            line_start_y = bar_y + bar_height  # バーの下から
            line_end_y = circle_y  # 円の上まで
            title_y = Inches(4.35)  # タイトル位置
            desc_y = Inches(4.60)   # 説明文位置
        
        # 縦線
        line_height = abs(line_end_y - line_start_y)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_x - Inches(0.015), min(line_start_y, line_end_y),
            Inches(0.03), line_height
        )
        line.fill.solid()
        line.fill.fore_color.rgb = timeline_colors[i % len(timeline_colors)]
        line.line.fill.background()
        
        # 丸いアイコン（円）
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            line_x - circle_size / 2, circle_y,
            circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = timeline_colors[i % len(timeline_colors)]
        circle.line.fill.background()
        
        # 円内のアイコン番号（完全中央配置）
        # 数字を円の中心に配置（上にずれるので下方向に大きく調整）
        num_box = slide.shapes.add_textbox(
            line_x - circle_size / 2, 
            circle_y + Inches(0.12),  # さらに下に移動
            circle_size, 
            circle_size - Inches(0.24)  # 高さを縮小
        )
        tf = num_box.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{i+1}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # テキストコンテンツ
        content_width = Inches(2.0)
        content_x = line_x - content_width / 2
        
        # タイトル
        title_box = slide.shapes.add_textbox(content_x, title_y, content_width, Inches(0.25))
        p = title_box.text_frame.paragraphs[0]
        p.text = item.get('title', '')
        p.font.name = theme.FONT_JP
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # 説明文
        desc_box = slide.shapes.add_textbox(content_x, desc_y, content_width, Inches(0.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item.get('description', '')
        p.font.name = theme.FONT_JP
        p.font.size = Pt(9)
        p.font.color.rgb = theme.TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
    
    return slide


# ============================================================
# テンプレートマップ
# ============================================================
TEMPLATES = {
    'cover': create_cover,
    'list': create_list,
    'card_grid': create_card_grid,
    'chart_left_text_right': create_chart_left_text_right,
    'text_left_chart_right': create_text_left_chart_right,
    'table_left_text_right': create_table_left_text_right,
    'text_left_image_right': create_text_left_image_right,
    'three_points_circle': create_three_points_circle,
    'point_list_compact': create_point_list_compact,
    'point_list_large': create_point_list_large,
    'timeline': create_timeline,
}


# ============================================================
# 3. データ定義 (Configuration) ★★★ ここを編集 ★★★
# ============================================================
# 資料を作成する際は、このセクションのみ編集してください。

# ------------------------------------------------------------
# 3-1. SLIDE_TEMPLATES: スライドテンプレート定義（再利用可能）
# ------------------------------------------------------------
# slide_id で参照。同じテンプレートを複数回使用可能。
SLIDE_TEMPLATES = {
    1: {
        "name": "表紙",
        "template": "cover",
        "default_data": {
            "title": "",
            "subtitle": "",
            "credit": ""
        }
    },
    2: {
        "name": "目次（リスト）",
        "template": "list",
        "default_data": {
            "title": "",
            "subtitle": "",
            "items": []
        }
    },
    3: {
        "name": "左グラフ・右テキスト（棒グラフ）",
        "template": "chart_left_text_right",
        "default_data": {
            "title": "",
            "subtitle": "",
            "chart_type": "COLUMN_CLUSTERED",
            "categories": [],
            "series": [],
            "text_items": []
        }
    },
    4: {
        "name": "左テキスト・右グラフ（円グラフ）",
        "template": "text_left_chart_right",
        "default_data": {
            "title": "",
            "subtitle": "",
            "chart_type": "PIE",
            "categories": [],
            "series": [],
            "text_items": []
        }
    },
    5: {
        "name": "左テーブル・右テキスト",
        "template": "table_left_text_right",
        "default_data": {
            "title": "",
            "subtitle": "",
            "columns": [],
            "rows": [],
            "text_items": []
        }
    },
    6: {
        "name": "左テキスト・右グラフ（折れ線）",
        "template": "text_left_chart_right",
        "default_data": {
            "title": "",
            "subtitle": "",
            "chart_type": "LINE",
            "categories": [],
            "series": [],
            "text_items": []
        }
    },
    7: {
        "name": "カードグリッド（4枚）",
        "template": "card_grid",
        "default_data": {
            "title": "",
            "subtitle": "",
            "items": []
        }
    },
    8: {
        "name": "左テキスト・右画像",
        "template": "text_left_image_right",
        "default_data": {
            "title": "",
            "subtitle": "",
            "text_items": []
        }
    },
    9: {
        "name": "3ポイント（丸囲み）",
        "template": "three_points_circle",
        "default_data": {
            "title": "",
            "subtitle": "",
            "items": []  # "見出し\n説明文" 形式
        }
    },
    10: {
        "name": "ポイントリスト（コンパクト）",
        "template": "point_list_compact",
        "default_data": {
            "title": "",
            "subtitle": "",
            "items": []  # "見出し\n説明文" 形式
        }
    },
    11: {
        "name": "ポイントリスト（大）",
        "template": "point_list_large",
        "default_data": {
            "items": []  # "見出し\n説明文" 形式（タイトルなし）
        }
    },
    12: {
        "name": "タイムライン",
        "template": "timeline",
        "default_data": {
            "title": "",
            "subtitle": "",
            "timeline_items": []  # {"year": "2024", "title": "...", "description": "..."} のリスト
        }
    }
}

# ------------------------------------------------------------
# 3-2. SLIDES_TO_USE: 使用するスライド番号を指定（順序反映、複数回使用OK）
# ------------------------------------------------------------
# 例: [1, 2, 3, 4, 5, 6, 7, 8, 1] → 表紙→目次→...→表紙（結び）
SLIDES_TO_USE = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 1]

# ------------------------------------------------------------
# 3-3. SLIDE_CONTENT: 各スライドのテキスト内容（output_index で個別設定）
# ------------------------------------------------------------
# output_index: 出力順（1始まり）。SLIDES_TO_USE の順序に対応。
# 指定がない項目はデフォルト値が使われる。
SLIDE_CONTENT = {
    # === 出力1枚目: 表紙 (slide_id=1) ===
    1: {
        "title": "2025年度 中期経営計画",
        "subtitle": "Sustainable Growth & Innovation",
        "credit": "2024.11.25 | 株式会社未来イノベーション"
    },
    
    # === 出力2枚目: 目次 (slide_id=2) ===
    2: {
        "title": "Agenda",
        "subtitle": "Structure",
        "items": [
            "市場環境分析と成長予測 (Market)",
            "2024年度 業績ハイライト (Review)",
            "2025年度 戦略フレームワーク (Strategy)",
            "デジタル変革と新サービス (DX)",
            "投資計画とリソース配分 (Resource)",
            "今後の展望と結論 (Conclusion)"
        ]
    },
    
    # === 出力3枚目: 市場分析・棒グラフ (slide_id=3) ===
    3: {
        "title": "Market Trends",
        "subtitle": "Analysis",
        "chart_type": "COLUMN_CLUSTERED",
        "categories": ["2023", "2024", "2025(E)", "2026(E)"],
        "series": [
            {"name": "IT Market", "values": [100, 102, 103, 104]},
            {"name": "DX Market", "values": [20, 35, 55, 80]}
        ],
        "text_items": [
            "DX市場は年率30%以上の高成長を維持",
            "既存IT市場は横ばい傾向が続く",
            "2025年が市場構造の転換点となる予測",
            "AI/クラウド領域への投資が急増中"
        ]
    },
    
    # === 出力4枚目: 市場シェア・円グラフ (slide_id=4) ===
    4: {
        "title": "Market Share",
        "subtitle": "Current Position",
        "chart_type": "PIE",
        "categories": ["当社", "A社", "B社", "その他"],
        "series": [
            {"name": "Share", "values": [35, 25, 20, 20]}
        ],
        "text_items": [
            "国内市場においてシェア35%で首位を維持",
            "2位A社との差は10ポイントに拡大",
            "「その他」に含まれる新興ベンチャーの台頭に注意",
            "来期目標はシェア40%の獲得"
        ]
    },
    
    # === 出力5枚目: 業績振り返り (slide_id=5) ===
    5: {
        "title": "FY2024 Review",
        "subtitle": "Performance",
        "columns": ["Metric", "Target", "Actual", "Rate"],
        "rows": [
            ["Sales", "100", "98", "98%"],
            ["Profit", "10", "12", "120%"],
            ["Clients", "500", "550", "110%"],
            ["Churn", "2.0%", "1.8%", "Good"]
        ],
        "text_items": [
            "利益面では目標を大きく超過達成 (120%)",
            "高収益なクラウドサービスの比率が向上",
            "売上高は一部案件の期ズレにより未達",
            "顧客基盤は順調に拡大中"
        ]
    },
    
    # === 出力6枚目: 戦略・折れ線グラフ (slide_id=6) ===
    6: {
        "title": "Sales Strategy",
        "subtitle": "Growth Plan",
        "chart_type": "LINE",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [
            {"name": "Cloud", "values": [20, 25, 32, 45]},
            {"name": "Consulting", "values": [15, 16, 16, 18]}
        ],
        "text_items": [
            "クラウド事業を成長ドライバーと位置付け",
            "コンサルティングは高付加価値化へシフト",
            "Q4に向けて新製品投入による加速を計画"
        ]
    },
    
    # === 出力7枚目: 新サービス・カードグリッド (slide_id=7) ===
    7: {
        "title": "New Solution Features",
        "subtitle": "Core Value",
        "items": [
            "All-in-One\n業務アプリを単一基盤に統合",
            "AI Native\n生成AIによる自動化を標準装備",
            "Security\n金融機関レベルの堅牢性",
            "UX/UI\n学習コストゼロのデザイン"
        ]
    },
    
    # === 出力8枚目: 今後の展望 (slide_id=8) ===
    8: {
        "title": "Future Vision",
        "subtitle": "Roadmap",
        "text_items": [
            "2025年: 国内シェアNo.1を獲得",
            "2026年: アジア主要都市へ拠点展開",
            "2027年: グローバルプラットフォームへ"
        ],
        "image_keyword": "future,technology"  # Unsplash検索キーワード
    },
    
    # === 出力9枚目: 3ポイント丸囲み (slide_id=9) ===
    9: {
        "title": "サービスの特長",
        "subtitle": "Features",
        "items": [
            "高速処理\nAI技術により従来の10倍の処理速度を実現",
            "安全性\n金融機関レベルのセキュリティ対応",
            "サポート\n24時間365日の専門サポート体制"
        ]
    },
    
    # === 出力10枚目: ポイントリストコンパクト (slide_id=10) ===
    10: {
        "title": "導入ステップ",
        "subtitle": "Process",
        "items": [
            "ヒアリング\n現状の課題と目標をヒアリングします",
            "提案\n最適なソリューションをご提案します",
            "導入\n最短2週間でサービス開始が可能です",
            "運用サポート\n導入後も継続的にサポートします"
        ]
    },
    
    # === 出力11枚目: ポイントリスト大 (slide_id=11) ===
    11: {
        "items": [
            "実績と信頼\n創業20年、累計1,000社以上の導入実績。業界トップクラスの顧客満足度を誇ります。",
            "柔軟なカスタマイズ\nお客様のニーズに合わせた柔軟なカスタマイズが可能。標準機能だけでなく独自要件にも対応します。",
            "継続的な進化\n毎月のアップデートで常に最新機能を提供。お客様のフィードバックを反映し続けます。"
        ]
    },
    
    # === 出力12枚目: タイムライン (slide_id=12) ===
    12: {
        "title": "ロードマップ",
        "subtitle": "Roadmap",
        "timeline_items": [
            {"year": "2024", "title": "Phase 1: 基盤構築", "description": "システム基盤の整備と初期導入を実施"},
            {"year": "2025", "title": "Phase 2: 展開", "description": "全社展開と機能拡張を推進"},
            {"year": "2026", "title": "Phase 3: 最適化", "description": "運用最適化とAI活用の本格化"},
            {"year": "2027", "title": "Phase 4: 革新", "description": "次世代プラットフォームへの進化"}
        ]
    },
    
    # === 出力13枚目: 結び・表紙再利用 (slide_id=1) ===
    13: {
        "title": "Next Stage",
        "subtitle": "Vision 2030",
        "credit": "株式会社未来イノベーション"
    }
}


# ============================================================
# 4. エンジン & 実行 (Engine & Execution) ※編集不要※
# ============================================================
# ※このセクションは触らなくてOKです。

def build_slide_data(output_index: int, slide_id: int) -> dict:
    """
    出力順(output_index)とスライドID(slide_id)から、
    実際のスライドデータを構築する。
    
    優先順位:
    1. SLIDE_CONTENT[output_index] の値
    2. SLIDE_TEMPLATES[slide_id]['default_data'] の値
    """
    template_info = SLIDE_TEMPLATES.get(slide_id, {})
    default_data = template_info.get('default_data', {}).copy()
    content_override = SLIDE_CONTENT.get(output_index, {})
    
    # default_data を content_override で上書きマージ
    merged_data = {**default_data, **content_override}
    
    return {
        "template": template_info.get('template', 'cover'),
        "template_name": template_info.get('name', '不明'),
        "data": merged_data
    }


def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 出力ファイル名はPythonファイル名と同じにする
    script_name = Path(__file__).stem  # 拡張子なしのファイル名
    output_path = Path(__file__).parent / "outputs" / f"{script_name}.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"🎨 テーマ: Modern Japan (Enhanced Charts Edition)")
    print(f"📊 スライド生成開始... (全{len(SLIDES_TO_USE)}枚)")
    print(f"📋 使用スライド: {SLIDES_TO_USE}")
    print("="*50)

    for output_index, slide_id in enumerate(SLIDES_TO_USE, start=1):
        # スライドデータを構築
        slide_info = build_slide_data(output_index, slide_id)
        template_name = slide_info['template']
        template_label = slide_info['template_name']
        data = slide_info['data']
        
        print(f"  [{output_index}] slide_id={slide_id} ({template_label})")
        print(f"      → title: {data.get('title', '(なし)')}")
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        if template_name in TEMPLATES:
            TEMPLATES[template_name](slide, data, CURRENT_THEME)
        else:
            print(f"      ⚠️  Warning: Template '{template_name}' not found.")

    prs.save(output_path)
    print("="*50)
    print(f"Generation completed: {output_path}")
    print("="*50)


def show_available_templates():
    """利用可能なテンプレート一覧を表示"""
    print("\n📖 利用可能なスライドテンプレート:")
    print("="*50)
    for slide_id, info in SLIDE_TEMPLATES.items():
        print(f"  {slide_id}: {info['name']} ({info['template']})")
    print("="*50)
    print("\n💡 使い方:")
    print("  1. SLIDES_TO_USE で使いたいスライド番号を指定")
    print("     例: [1, 2, 3, 4, 1]  # 表紙→目次→...→表紙")
    print("  2. SLIDE_CONTENT で各スライドのテキストを編集")
    print("     キー = 出力順（1始まり）")
    print()


if __name__ == "__main__":
    show_available_templates()
    generate_presentation()
