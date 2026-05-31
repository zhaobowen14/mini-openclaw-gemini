from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


class Theme:
    PRIMARY = RGBColor(0, 82, 155)
    SECONDARY = RGBColor(0, 114, 198)
    TEXT = RGBColor(51, 51, 51)
    TEXT_LIGHT = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)
    BACKGROUND = RGBColor(255, 255, 255)
    CHART_COLORS = [
        RGBColor(0, 82, 155),
        RGBColor(0, 114, 198),
        RGBColor(65, 143, 222),
    ]
    FONT_JP = "メイリオ"
    FONT_EN = "Segoe UI"


THEME = Theme()


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME.BACKGROUND


def add_left_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME.PRIMARY
    bar.line.fill.background()


def add_header(slide, title, subtitle=""):
    set_background(slide)
    add_left_bar(slide)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.5), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME.PRIMARY
    line.line.fill.background()


def create_cover(slide, data):
    set_background(slide)
    add_left_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.5), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY

    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY

    credit = data.get("credit", "")
    if credit:
        credit_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(8.5), Inches(0.3))
        p = credit_box.text_frame.paragraphs[0]
        p.text = credit
        p.font.name = THEME.FONT_JP
        p.font.size = Pt(12)
        p.font.color.rgb = THEME.TEXT_LIGHT


def create_list(slide, data):
    add_header(slide, data.get("title", ""), data.get("subtitle", ""))
    items = data.get("items", [])

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = THEME.FONT_JP
        p.font.size = Pt(18)
        p.font.color.rgb = THEME.TEXT
        p.space_after = Pt(12)


def create_chart(slide, data):
    add_header(slide, data.get("title", ""), data.get("subtitle", ""))

    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for s in data.get("series", []):
        chart_data.add_series(s["name"], s["values"])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.7), Inches(1.5), Inches(8.0), Inches(3.5),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = THEME.FONT_JP

    for i, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = THEME.CHART_COLORS[i % len(THEME.CHART_COLORS)]


TEMPLATES = {
    "cover": create_cover,
    "list": create_list,
    "chart": create_chart,
}


SLIDES_TO_USE = [1, 2, 3]

SLIDE_CONTENT = {
    1: {
        "template": "cover",
        "title": "2025年度 中期経営計画",
        "subtitle": "Sustainable Growth & Innovation",
        "credit": "2024.11.25 | 株式会社未来イノベーション",
    },
    2: {
        "template": "list",
        "title": "Agenda",
        "subtitle": "Structure",
        "items": [
            "市場環境分析",
            "業績ハイライト",
            "戦略フレームワーク",
            "今後の展望",
        ],
    },
    3: {
        "template": "chart",
        "title": "Market Trends",
        "subtitle": "Analysis",
        "categories": ["2023", "2024", "2025(E)", "2026(E)"],
        "series": [
            {"name": "IT Market", "values": [100, 102, 103, 104]},
            {"name": "DX Market", "values": [20, 35, 55, 80]},
        ],
    },
}


def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_no in SLIDES_TO_USE:
        data = SLIDE_CONTENT[slide_no]
        template_name = data["template"]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        TEMPLATES[template_name](slide, data)

    output_dir = Path("./")
    output_path = output_dir / "minimal_presentation.pptx"
    prs.save(output_path)
    print(f"Generation completed: {output_path}")


if __name__ == "__main__":
    generate_presentation()