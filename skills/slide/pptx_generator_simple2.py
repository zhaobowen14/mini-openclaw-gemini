from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls




class Theme:
    # Primary: コーポレートブルー（濃い青）
    PRIMARY = RGBColor(0, 82, 155)  # #00529B
    
    # Secondary: アクセントブルー（明るい青）
    SECONDARY = RGBColor(0, 114, 198)  # #0072C6
    
    # Accent: ポイントカラー（さらに明るい青）
    ACCENT = RGBColor(65, 143, 222)  # #418FDE
    
    # Chart Colors (Blue Gradation Palette)
    CHART_COLORS = [
        RGBColor(0, 82, 155),    # Primary Blue (濃)
        RGBColor(0, 114, 198),   # Secondary Blue
        RGBColor(65, 143, 222),  # Accent Blue
        RGBColor(130, 180, 230), # Light Blue
        RGBColor(180, 210, 240), # Lighter Blue
        RGBColor(220, 235, 250)  # Very Light Blue
    ]
    
    # Base Colors
    BACKGROUND = RGBColor(255, 255, 255)  # 白
    TEXT = RGBColor(51, 51, 51)           # 濃いグレー #333333
    TEXT_LIGHT = RGBColor(102, 102, 102)  # ミディアムグレー #666666
    WHITE = RGBColor(255, 255, 255)
    
    # UI Elements
    DIVIDER_LINE = RGBColor(0, 82, 155)   # 青のライン
    BOX_BORDER = RGBColor(200, 220, 240)  # 薄い青のボーダー
    LIGHT_BG = RGBColor(240, 247, 255)    # 薄い青の背景
    
    # Number/Point Colors
    POINT_NUMBER = RGBColor(0, 82, 155)   # ポイント番号の青
    
    # Fonts
    FONT_JP = "メイリオ"
    FONT_EN = "Segoe UI"


THEME = Theme()


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME.BACKGROUND


def add_left_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME.PRIMARY
    bar.line.fill.background()


def add_header(slide, title, subtitle=""):
    set_background(slide)
    add_left_bar(slide)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.5), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = THEME.PRIMARY
    line.line.fill.background()


def create_cover(slide, data):
    set_background(slide)
    add_left_bar(slide)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.5), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY

    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY


def create_list(slide, data):
    add_header(slide, data.get("title", ""), data.get("subtitle", ""))
    items = data.get("items", [])

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = THEME.FONT_JP
        p.font.size = Pt(18)
        p.font.color.rgb = THEME.TEXT
        p.space_after = Pt(12)


def create_chart(slide, data):
    add_header(slide, data.get("title", ""), data.get("subtitle", ""))

    chart_data = CategoryChartData()
    chart_data.categories = data.get("categories", [])
    for s in data.get("series", []):
        chart_data.add_series(s["name"], s["values"])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.7), Inches(1.5), Inches(8.0), Inches(3.5),
        chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.name = THEME.FONT_JP

    for i, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = THEME.CHART_COLORS[i % len(THEME.CHART_COLORS)]



# --- Helper Functions ---


def _enable_bullet(paragraph):
    """段落に対してPowerPoint本来の箇条書き属性を付与"""
    pPr = paragraph._p.get_or_add_pPr()
    buChar = parse_xml('<a:buChar %s char="•"/>' % nsdecls('a'))
    if pPr.find(f"{{{nsdecls('a')}}}buChar") is None:
        pPr.insert(0, buChar)
    pPr.set('marL', '457200')
    pPr.set('indent', '-228600')


def _create_text_list_bullet(slide, items, theme, x, y, w):
    """バレット付きテキストリスト - シンプル箇条書き"""
    tb = slide.shapes.add_textbox(x, y, w, Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.color.rgb = theme.TEXT
        p.space_after = Pt(10)
        _enable_bullet(p)




def _create_table_at(slide, data, theme, x, y, w):
    """テーブル - 青ヘッダースタイル"""
    columns = data.get('columns', [])
    rows = data.get('rows', [])
    rows_count = len(rows) + 1
    cols_count = len(columns)
    
    row_height = Inches(0.4)
    table_height = row_height * rows_count
    
    shape = slide.shapes.add_table(rows_count, cols_count, x, y, w, table_height)
    table = shape.table
    
    # ヘッダー行（青背景・白文字）
    for i, col in enumerate(columns):
        cell = table.cell(0, i)
        cell.text = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.PRIMARY  # 青背景
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = theme.WHITE  # 白文字
        p.alignment = PP_ALIGN.CENTER

    # データ行
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            # 交互に薄い青背景
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.LIGHT_BG
            else:
                cell.fill.background()
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = theme.TEXT
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

def _setup_header(slide, data, theme):
    """ヘッダー - 青ライン＋タイトル"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = theme.BACKGROUND

    # 左側の青いバー（細め）
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = theme.PRIMARY
    left_bar.line.fill.background()

    # タイトル下の青いライン
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = theme.PRIMARY
    line.line.fill.background()

    if 'subtitle' in data:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = data['subtitle']
        p.font.name = theme.FONT_EN
        p.font.size = Pt(10)
        p.font.color.rgb = theme.SECONDARY
        p.font.bold = True

    title = data.get('title', '')
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(8), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = theme.FONT_JP
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = theme.PRIMARY  # タイトルを青に




def create_table_left_text_right(slide, data, theme=THEME):
    """左：テーブル / 右：テキスト"""
    _setup_header(slide, data, theme)
    _create_table_at(slide, data, theme, Inches(0.5), Inches(1.4), Inches(5.5))
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(6.2), Inches(1.4), Inches(3.3))


def create_three_points_circle(slide, data, theme=THEME):
    """3ポイント - 丸囲み番号スタイル"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    
    start_y = Inches(1.5)
    col_width = Inches(2.8)
    gap = Inches(0.4)
    start_x = Inches(0.7)
    
    for i, item in enumerate(items[:3]):
        x = start_x + i * (col_width + gap)
        
        # 丸い枠（円）
        circle_x = x + Inches(0.8)
        circle_y = start_y
        circle_size = Inches(1.4)
        
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, circle_size, circle_size)
        circle.fill.background()
        circle.line.color.rgb = theme.PRIMARY
        circle.line.width = Pt(2)
        
        # 番号（円の中央 - 手動で位置調整）
        # 円の中心: circle_y + circle_size/2 = start_y + 0.7インチ
        # テキスト高さ約0.5インチとして、中央に配置
        text_height = Inches(0.55)
        text_y = circle_y + (circle_size - text_height) / 2 + Inches(0.02)  # 微調整
        
        num_box = slide.shapes.add_textbox(circle_x, text_y, circle_size, text_height)
        tf = num_box.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        tf.word_wrap = False
        
        p = tf.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        p.alignment = PP_ALIGN.CENTER  # 水平中央揃え
        
        # 見出しと説明文を分離
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        # 見出し（円の下）
        head_box = slide.shapes.add_textbox(x, start_y + Inches(1.55), col_width, Inches(0.4))
        head_box.text_frame.word_wrap = True
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # 説明文
        if desc:
            desc_box = slide.shapes.add_textbox(x, start_y + Inches(1.95), col_width, Inches(1.5))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(10)
            p.font.color.rgb = theme.TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER






def _create_chart_at(slide, data, theme, x, y, w, h):
    """グラフ作成（色設定・円グラフ対応）"""
    
    chart_type_str = data.get('chart_type', 'COLUMN_CLUSTERED')
    
    # 円グラフ（Pie Chart）の場合のデータ構造は異なる
    if chart_type_str == 'PIE':
        chart_data = ChartData()
        chart_data.categories = data.get('categories', [])
        # 円グラフは最初のシリーズのみ使用
        series_list = data.get('series', [])
        if series_list:
            chart_data.add_series(series_list[0]['name'], series_list[0]['values'])
        chart_type = XL_CHART_TYPE.PIE
    else:
        chart_data = CategoryChartData()
        chart_data.categories = data.get('categories', [])
        series_list = data.get('series', [])
        for s in series_list:
            chart_data.add_series(s['name'], s['values'])
        chart_type = getattr(XL_CHART_TYPE, chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # チャート追加
    chart = slide.shapes.add_chart(chart_type, x, y, w, h, chart_data).chart
    
    # スタイル・凡例設定
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = theme.FONT_JP
    chart.chart_title.has_text_frame = False
    
    # データラベル（円グラフ用）
    if chart_type == XL_CHART_TYPE.PIE:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(9)
        data_labels.font.color.rgb = theme.TEXT
    
    # シリーズの色設定
    # python-pptxのchart.seriesのAPIは限定的だが、可能な範囲で適用
    try:
        if chart_type == XL_CHART_TYPE.PIE:
            # 円グラフは「ポイント」ごとに色を変える
            plot = chart.plots[0]
            series = plot.series[0]
            for i, point in enumerate(series.points):
                fill = point.format.fill
                fill.solid()
                # 用意したパレットを循環
                color = theme.CHART_COLORS[i % len(theme.CHART_COLORS)]
                fill.fore_color.rgb = color
        else:
            # 棒・折れ線は「シリーズ」ごとに色を変える
            for i, series in enumerate(chart.series):
                fill = series.format.fill
                fill.solid()
                color = theme.CHART_COLORS[i % len(theme.CHART_COLORS)]
                fill.fore_color.rgb = color
                
                # 折れ線グラフの場合、線自体の色も変える必要があるが
                # line.colorプロパティへのアクセスはバージョンにより異なる場合あり
                if hasattr(series.format.line, 'color'):
                     series.format.line.color.rgb = color

    except Exception as e:
        # 色設定に失敗してもグラフ自体は生成されるようにする
        print(f"Note: Color formatting skipped for some elements ({e})")



def create_text_left_chart_right(slide, data, theme=THEME):
    """左：テキスト / 右：グラフ"""
    _setup_header(slide, data, theme)
    _create_text_list_bullet(slide, data.get('text_items', []), theme, Inches(0.5), Inches(1.4), Inches(3.5))
    _create_chart_at(slide, data, theme, Inches(4.2), Inches(1.4), Inches(5.3), Inches(3.8))

def create_point_list_compact(slide, data, theme=THEME):
    """ポイントリスト - コンパクト（番号 + 見出し + 説明）"""
    _setup_header(slide, data, theme)
    items = data.get('items', [])
    
    start_y = Inches(1.4)
    item_height = Inches(0.85)
    
    for i, item in enumerate(items[:4]):
        y_pos = start_y + i * item_height
        
        # 番号
        num_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(0.5), Inches(0.5))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        
        # 見出しと説明文を分離
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        # 見出し
        head_box = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(8.0), Inches(0.35))
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        
        # 説明文
        if desc:
            desc_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.35), Inches(8.0), Inches(0.45))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(11)
            p.font.color.rgb = theme.TEXT_LIGHT


TEMPLATES = {
    "cover": create_cover,
    "list": create_list,
    "chart": create_chart,
    'table_left_text_right': create_table_left_text_right,
    'text_left_chart_right': create_text_left_chart_right,
    'three_points_circle': create_three_points_circle,
    'point_list_compact': create_point_list_compact,
}

SLIDES_TO_USE = [1, 2, 5, 9, 10]

SLIDE_CONTENT = {
    1: {
        "template": "cover",
        "title": "2025年度 中期経営計画",
        "subtitle": "Sustainable Growth & Innovation",
    },
    2: {
        "template": "list",
        "title": "Agenda",
        "subtitle": "Structure",
        "items": [
            "市場環境分析",
            "業績ハイライト",
            "戦略フレームワーク",
            "今後の展望",
        ],
    },

    5: {
        "template": "table_left_text_right",
        "title": "FY2024 Review",
        "subtitle": "Performance",
        "columns": ["Metric", "Target", "Actual", "Rate"],
        "rows": [
            ["Sales", "100", "98", "98%"],
            ["Profit", "10", "12", "120%"],
            ["Clients", "500", "550", "110%"],
            ["Churn", "2.0%", "1.8%", "Good"]
        ],
        "text_items": [
            "利益面では目標を大きく超過達成 (120%)",
            "高収益なクラウドサービスの比率が向上",
            "売上高は一部案件の期ズレにより未達",
            "顧客基盤は順調に拡大中"
        ]
    },
    
    # === 折れ線グラフ ===
    6: {
        "template": "text_left_chart_right",
        "title": "Sales Strategy",
        "subtitle": "Growth Plan",
        "chart_type": "LINE",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [
            {"name": "Cloud", "values": [20, 25, 32, 45]},
            {"name": "Consulting", "values": [15, 16, 16, 18]}
        ],
        "text_items": [
            "クラウド事業を成長ドライバーと位置付け",
            "コンサルティングは高付加価値化へシフト",
            "Q4に向けて新製品投入による加速を計画"
        ]
    },

        # === 出力9枚目: 3ポイント丸囲み (slide_id=9) ===
    9: {
        "template": "three_points_circle",
        "title": "サービスの特長",
        "subtitle": "Features",
        "items": [
            "高速処理\nAI技術により従来の10倍の処理速度を実現",
            "安全性\n金融機関レベルのセキュリティ対応",
            "サポート\n24時間365日の専門サポート体制"
        ]
    },
    
    # === 出力10枚目: ポイントリストコンパクト (slide_id=10) ===
    10: {
        "template": "point_list_compact",
        "title": "導入ステップ",
        "subtitle": "Process",
        "items": [
            "ヒアリング\n現状の課題と目標をヒアリングします",
            "提案\n最適なソリューションをご提案します",
            "導入\n最短2週間でサービス開始が可能です",
            "運用サポート\n導入後も継続的にサポートします"
        ]
    },
    
}


def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for slide_no in SLIDES_TO_USE:
        data = SLIDE_CONTENT[slide_no]
        template_name = data["template"]

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        TEMPLATES[template_name](slide, data)

    output_dir = Path("./")
    output_path = output_dir / "minimal_presentation.pptx"
    prs.save(output_path)
    print(f"Generation completed: {output_path}")


if __name__ == "__main__":
    generate_presentation()