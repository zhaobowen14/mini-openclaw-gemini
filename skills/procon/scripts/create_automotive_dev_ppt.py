# =========================================================
# 参考コード
# - 1スライド1メッセージ
# - メリデメ表、スコアリング、時間軸を含む
# =========================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pathlib import Path



OUT_PATH = Path('/mnt/data/automotive_bev_decision_deck_default_tables.pptx')

# Slide size: 16:9
SLIDE_W = 13.333
SLIDE_H = 7.5

# Colors
NAVY = RGBColor(22, 38, 66)
BLUE = RGBColor(36, 96, 168)
LIGHT_BLUE = RGBColor(226, 237, 251)
ORANGE = RGBColor(236, 142, 46)
LIGHT_ORANGE = RGBColor(255, 238, 214)
RED = RGBColor(206, 69, 62)
LIGHT_RED = RGBColor(252, 230, 230)
GREEN = RGBColor(47, 126, 92)
LIGHT_GREEN = RGBColor(225, 243, 234)
GRAY = RGBColor(84, 96, 112)
LIGHT_GRAY = RGBColor(244, 246, 249)
MID_GRAY = RGBColor(215, 220, 228)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(20, 24, 30)

# Default table style
# 表はPowerPointのデフォルト表に寄せるため、セル背景色・罫線色は指定しない。
TABLE_TEXT = BLACK
TABLE_HEADER_TEXT = BLACK

FONT_JP = 'Meiryo'
FONT_EN = 'Arial'

# Tracking for layout validation
_slide_elements = {}


def emu_to_inch(x):
    return x / 914400


def register(slide, name, x, y, w, h, ignore_overlap=False):
    slide_id = id(slide)
    _slide_elements.setdefault(slide_id, []).append({
        'name': name, 'x': x, 'y': y, 'w': w, 'h': h, 'ignore': ignore_overlap
    })


def rects_overlap(a, b, tol=0.02):
    return not (
        a['x'] + a['w'] <= b['x'] + tol or
        b['x'] + b['w'] <= a['x'] + tol or
        a['y'] + a['h'] <= b['y'] + tol or
        b['y'] + b['h'] <= a['y'] + tol
    )


def warnIfSlideHasOverlaps(slide):
    """警告用: 意図しない要素重なりを検出する。"""
    elems = [e for e in _slide_elements.get(id(slide), []) if not e['ignore']]
    warnings = []
    for i in range(len(elems)):
        for j in range(i + 1, len(elems)):
            if rects_overlap(elems[i], elems[j]):
                warnings.append(f"Overlap: {elems[i]['name']} vs {elems[j]['name']}")
    return warnings


def warnIfSlideElementsOutOfBounds(slide):
    """警告用: スライド外にはみ出した要素を検出する。"""
    warnings = []
    for e in _slide_elements.get(id(slide), []):
        if e['x'] < -0.01 or e['y'] < -0.01 or e['x'] + e['w'] > SLIDE_W + 0.01 or e['y'] + e['h'] > SLIDE_H + 0.01:
            warnings.append(f"OutOfBounds: {e['name']} {e}")
    return warnings


def set_text_frame_style(tf, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font=FONT_JP):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.bold = bold


def add_text(slide, text, x, y, w, h, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, name='text'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.name = name
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    register(slide, name, x, y, w, h)
    return tb


def add_shape_text(slide, text, x, y, w, h, fill=LIGHT_GRAY, line=MID_GRAY, font_size=16, color=BLACK, bold=False, align=PP_ALIGN.CENTER, radius=True, name='shape'):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.name = name
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_JP
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    register(slide, name, x, y, w, h)
    return shp


def add_header(slide, title, subtitle=None):
    add_text(slide, title, 0.55, 0.32, 9.9, 0.48, font_size=25, color=NAVY, bold=True, name='slide_title')
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.88, 9.8, 0.28, font_size=12.5, color=GRAY, name='slide_subtitle')
    # small page marker line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.20), Inches(12.2), Inches(0.02))
    line.name = 'header_line'
    line.fill.solid(); line.fill.fore_color.rgb = MID_GRAY
    line.line.fill.background()
    register(slide, 'header_line', 0.55, 1.20, 12.2, 0.02, ignore_overlap=True)


def add_bullets_to_cell(cell, bullets, font_size=10.5, color=BLACK):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = FONT_JP
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(0)
        p.space_before = Pt(0)


def set_cell(cell, text=None, fill=None, color=BLACK, bold=False, font_size=11, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    if text is not None:
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT_JP
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold



def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # -----------------------------
    # Slide 1: Title
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = LIGHT_GRAY
    add_text(slide, '自動車開発における\nBEV導入メリット・デメリット分析', 0.75, 1.25, 8.2, 1.25, font_size=34, color=NAVY, bold=True, name='title')
    add_text(slide, 'PowerPoint直接出力版｜技術選定・経営判断向け', 0.80, 2.72, 7.5, 0.42, font_size=17, color=GRAY, name='subtitle')
    add_shape_text(slide, '結論\nBEVへ段階的に投資しつつ、短期収益はHEVで確保', 9.20, 1.25, 3.10, 1.55, fill=WHITE, line=MID_GRAY, font_size=17, color=NAVY, bold=True, name='decision_box')
    add_shape_text(slide, '比較対象\nBEV / HEV / ICE', 0.82, 4.1, 2.95, 1.10, fill=LIGHT_BLUE, line=LIGHT_BLUE, font_size=18, color=NAVY, bold=True, name='scope_box')
    add_shape_text(slide, '評価軸\n性能・コスト・規制・信頼性', 4.02, 4.1, 3.65, 1.10, fill=LIGHT_GREEN, line=LIGHT_GREEN, font_size=18, color=NAVY, bold=True, name='criteria_box')
    add_shape_text(slide, '出力\n比較表・スコア・時間軸', 7.92, 4.1, 3.65, 1.10, fill=LIGHT_ORANGE, line=LIGHT_ORANGE, font_size=18, color=NAVY, bold=True, name='output_box')
    add_text(slide, '※数値は意思決定用の整理例。実案件では自社データ・市場データで更新してください。', 0.82, 6.78, 9.0, 0.28, font_size=10.5, color=GRAY, name='note')

    # -----------------------------
    # Slide 2: Executive Summary
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    add_header(slide, '意思決定サマリー', 'BEVは規制対応と効率で優位。一方で、短期は電池コストと供給リスクが制約。')
    add_shape_text(slide, '推奨方針', 0.65, 1.55, 2.15, 0.45, fill=NAVY, line=NAVY, font_size=16, color=WHITE, bold=True, name='label_recommend')
    add_shape_text(slide, 'BEV専用プラットフォームを中期投資テーマとし、短期はHEVで収益と量販を維持する。', 0.65, 2.10, 5.75, 1.40, fill=LIGHT_BLUE, line=LIGHT_BLUE, font_size=20, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='recommend_card')
    add_shape_text(slide, '主要リスク', 6.75, 1.55, 2.15, 0.45, fill=RED, line=RED, font_size=16, color=WHITE, bold=True, name='label_risk')
    add_shape_text(slide, '電池コスト、原材料供給、寒冷地性能、充電インフラの4点がBEV導入判断のボトルネック。', 6.75, 2.10, 5.75, 1.40, fill=LIGHT_RED, line=LIGHT_RED, font_size=20, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='risk_card')
    add_shape_text(slide, '実行優先順位', 0.65, 4.08, 2.15, 0.45, fill=GREEN, line=GREEN, font_size=16, color=WHITE, bold=True, name='label_action')
    add_shape_text(slide, '1. 高価格帯BEVから投入\n2. 電池調達とBMS開発を強化\n3. HEVで短期CFを維持', 0.65, 4.62, 5.75, 1.58, fill=LIGHT_GREEN, line=LIGHT_GREEN, font_size=19, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='action_card')
    add_shape_text(slide, '判断の見方', 6.75, 4.08, 2.15, 0.45, fill=ORANGE, line=ORANGE, font_size=16, color=WHITE, bold=True, name='label_view')
    add_shape_text(slide, 'BEVは長期戦略で優位、HEVは短期収益で優位。ICEは規制対応コスト増に注意。', 6.75, 4.62, 5.75, 1.58, fill=LIGHT_ORANGE, line=LIGHT_ORANGE, font_size=19, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='view_card')

    # -----------------------------
    # Slide 3: Pros/Cons table - BEV vs HEV
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    add_header(slide, 'メリット・デメリット比較: BEV vs HEV', '実務判断では、BEVは長期規制対応、HEVは短期収益性で評価する。')
    rows, cols = 3, 4
    x, y, w, h = 0.45, 1.50, 12.45, 5.10
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table_shape.name = 'bev_hev_table'
    register(slide, 'bev_hev_table', x, y, w, h)
    table = table_shape.table
    widths = [1.50, 4.20, 4.65, 2.10]
    for i, width in enumerate(widths):
        table.columns[i].width = Inches(width)
    heights = [0.55, 2.25, 2.30]
    for i, height in enumerate(heights):
        table.rows[i].height = Inches(height)
    headers = ['案', 'メリット', 'デメリット', '判断']
    for c, head in enumerate(headers):
        set_cell(table.cell(0, c), head, fill=None, color=TABLE_HEADER_TEXT, bold=False, font_size=13.5, align=PP_ALIGN.CENTER)
    data = [
        ['BEV',
         ['CO2排出 0g/kmでZEV規制に対応', '効率 約90%でエネルギー利用に優位', '部品点数 -30〜40%で構造を簡素化', '低重心で操縦安定性を高めやすい'],
         ['電池コストが車両価格の30〜40%', '急速充電インフラに依存', '寒冷地で航続距離 -30〜40%', 'リチウム等の供給不安が残る'],
         '高\n長期投資の主軸'],
        ['HEV',
         ['既存インフラをそのまま活用可能', 'ICE比で燃費 +30〜50%', '航続距離不安が小さい', '短期の排出量低減に有効'],
         ['エンジンとモーターで構造が複雑', 'ECU統合難易度が高い', '開発工数 +20%を想定', '長期ZEV規制では不利'],
         '中\n短期収益の柱']
    ]
    for r, row in enumerate(data, start=1):
        option, pros, cons, impact = row
        set_cell(table.cell(r, 0), option, fill=None, color=TABLE_TEXT, bold=False, font_size=17, align=PP_ALIGN.CENTER)
        add_bullets_to_cell(table.cell(r, 1), pros, font_size=12.2)
        add_bullets_to_cell(table.cell(r, 2), cons, font_size=12.2)
        set_cell(table.cell(r, 3), impact, fill=None, color=TABLE_TEXT, bold=False, font_size=14, align=PP_ALIGN.CENTER)
    add_text(slide, 'ポイント: BEVは「規制・効率」、HEVは「コスト・利便性」で優位。二者択一ではなく段階的移行が現実的。', 0.55, 6.82, 11.6, 0.28, font_size=12.0, color=GRAY, name='pros_cons_note')

    # -----------------------------
    # Slide 4: ICE positioning
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    add_header(slide, 'ICEの位置づけと長期リスク', 'ICEは短期では低コスト・高信頼だが、長期は規制対応コストが増える。')
    add_shape_text(slide, '強み', 0.70, 1.55, 1.60, 0.45, fill=GREEN, line=GREEN, font_size=16, color=WHITE, bold=True, name='ice_pro_label')
    add_shape_text(slide, '技術成熟により品質が安定\n既存生産ラインを活用可能\n寒冷地・長距離で性能が安定', 0.70, 2.12, 3.80, 2.15, fill=LIGHT_GREEN, line=LIGHT_GREEN, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='ice_pro_card')
    add_shape_text(slide, '弱み', 4.82, 1.55, 1.60, 0.45, fill=RED, line=RED, font_size=16, color=WHITE, bold=True, name='ice_con_label')
    add_shape_text(slide, '排出ガス規制対応コスト増\n熱効率に構造的な限界\n将来的な販売禁止・都市部規制', 4.82, 2.12, 3.80, 2.15, fill=LIGHT_RED, line=LIGHT_RED, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='ice_con_card')
    add_shape_text(slide, '実務判断', 8.94, 1.55, 1.90, 0.45, fill=NAVY, line=NAVY, font_size=16, color=WHITE, bold=True, name='ice_judgement_label')
    add_shape_text(slide, '新規投資の主軸にはしにくい。\n既存ラインの回収と地域別需要への対応として限定的に活用する。', 8.94, 2.12, 3.40, 2.15, fill=LIGHT_BLUE, line=LIGHT_BLUE, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.LEFT, name='ice_judgement_card')
    add_shape_text(slide, 'ICEは「撤退」ではなく「縮小管理」: 収益回収、地域最適化、規制対応費の上限管理が重要', 1.05, 5.45, 11.20, 0.85, fill=NAVY, line=NAVY, font_size=20, color=WHITE, bold=True, name='ice_takeaway')

    # -----------------------------
    # Slide 4: Scoring
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    add_header(slide, '評価軸別スコアリング', 'スコアは0〜100、重み合計は100%。BEVは総合81、HEVは78。')
    chart_data = CategoryChartData()
    chart_data.categories = ['性能', 'コスト', '規制対応', '信頼性']
    chart_data.add_series('BEV', (90, 60, 100, 80))
    chart_data.add_series('HEV', (75, 85, 65, 90))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.70), Inches(1.65), Inches(6.60), Inches(4.70), chart_data).chart
    register(slide, 'score_chart', 0.70, 1.65, 6.60, 4.70)
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.maximum_scale = 110
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 20
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.chart_title.text_frame.text = ''
    for series in chart.series:
        series.has_data_labels = True
        series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        series.data_labels.font.size = Pt(9)
    # Summary table
    rows, cols = 6, 4
    x, y, w, h = 7.70, 1.70, 4.80, 4.55
    tshape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tshape.name = 'score_table'
    register(slide, 'score_table', x, y, w, h)
    t = tshape.table
    for i, width in enumerate([1.40, 0.85, 0.80, 1.75]):
        t.columns[i].width = Inches(width)
    for i, height in enumerate([0.48, 0.72, 0.72, 0.72, 0.72, 0.72]):
        t.rows[i].height = Inches(height)
    header = ['評価軸', '重み', 'BEV', 'HEV']
    for c, htxt in enumerate(header):
        set_cell(t.cell(0, c), htxt, fill=None, color=TABLE_HEADER_TEXT, bold=False, font_size=11.5, align=PP_ALIGN.CENTER)
    score_rows = [
        ['性能', '25%', '90', '75'],
        ['コスト', '30%', '60', '85'],
        ['規制対応', '25%', '100', '65'],
        ['信頼性', '20%', '80', '90'],
        ['合計', '100%', '81', '78']
    ]
    for r, row in enumerate(score_rows, start=1):
        for c, txt in enumerate(row):
            set_cell(t.cell(r, c), txt, fill=None, color=TABLE_TEXT, bold=False, font_size=11.3, align=PP_ALIGN.CENTER)
    add_shape_text(slide, '読み取り\nBEVは規制・性能で優位。HEVはコスト・信頼性で優位。', 0.70, 6.50, 11.80, 0.55, fill=LIGHT_BLUE, line=LIGHT_BLUE, font_size=17, color=NAVY, bold=True, name='score_takeaway')

    # -----------------------------
    # Slide 5: Roadmap
    # -----------------------------
    slide = prs.slides.add_slide(blank)
    add_header(slide, '実行ロードマップ', '短期はHEVで収益を維持し、中期以降にBEV投資を加速する。')
    # Horizontal timeline base
    add_shape_text(slide, '短期\n〜3年', 0.75, 1.70, 3.55, 0.72, fill=LIGHT_GREEN, line=GREEN, font_size=20, color=NAVY, bold=True, name='short_label')
    add_shape_text(slide, '中期\n3〜7年', 4.85, 1.70, 3.55, 0.72, fill=LIGHT_ORANGE, line=ORANGE, font_size=20, color=NAVY, bold=True, name='mid_label')
    add_shape_text(slide, '長期\n7年以上', 8.95, 1.70, 3.55, 0.72, fill=LIGHT_RED, line=RED, font_size=20, color=NAVY, bold=True, name='long_label')
    add_shape_text(slide, 'HEV販売比率を維持\n高価格帯BEVを限定投入\n電池サプライヤーを複線化', 0.75, 2.78, 3.55, 2.25, fill=WHITE, line=MID_GRAY, font_size=17, color=BLACK, bold=False, align=PP_ALIGN.LEFT, name='short_card')
    add_shape_text(slide, '次世代電池を評価\nBEV専用PFへ移行\nBMS・熱管理を内製強化', 4.85, 2.78, 3.55, 2.25, fill=WHITE, line=MID_GRAY, font_size=17, color=BLACK, bold=False, align=PP_ALIGN.LEFT, name='mid_card')
    add_shape_text(slide, 'ZEV規制に完全適合\nリサイクル網を構築\nLCC最適化で資源リスク低減', 8.95, 2.78, 3.55, 2.25, fill=WHITE, line=MID_GRAY, font_size=17, color=BLACK, bold=False, align=PP_ALIGN.LEFT, name='long_card')
    add_shape_text(slide, '次のアクション: 自社のBOM、地域別規制、販売台数シナリオを入れてスコアを再計算する', 1.15, 6.12, 11.10, 0.70, fill=NAVY, line=NAVY, font_size=19, color=WHITE, bold=True, name='next_action')

    # Validation
    all_warnings = []
    for i, s in enumerate(prs.slides, start=1):
        for wmsg in warnIfSlideHasOverlaps(s):
            all_warnings.append(f'Slide {i}: {wmsg}')
        for wmsg in warnIfSlideElementsOutOfBounds(s):
            all_warnings.append(f'Slide {i}: {wmsg}')
    # Keep non-fatal: intentionally no overlaps should be detected due to careful placement.
    if all_warnings:
        print('Layout warnings:')
        for msg in all_warnings:
            print(' -', msg)
    else:
        print('Layout validation passed: no unintended overlaps or out-of-bounds elements.')

    prs.save(OUT_PATH)
    print(f'Saved: {OUT_PATH}')


if __name__ == '__main__':
    build_deck()
