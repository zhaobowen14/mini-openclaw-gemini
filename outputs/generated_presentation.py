# quantum_presentation.py
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

class Theme:
    PRIMARY = RGBColor(0, 82, 155)
    SECONDARY = RGBColor(0, 114, 198)
    ACCENT = RGBColor(65, 143, 222)
    BACKGROUND = RGBColor(255, 255, 255)
    TEXT = RGBColor(51, 51, 51)
    TEXT_LIGHT = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)
    FONT_JP = "メイリオ"
    FONT_EN = "Segoe UI"

THEME = Theme()

def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME.BACKGROUND

def add_left_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME.PRIMARY
    bar.line.fill.background()

def add_header(slide, title, subtitle=""):
    set_background(slide)
    add_left_bar(slide)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.3))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(8.5), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(9.0), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = THEME.PRIMARY
    line.line.fill.background()

def create_cover(slide, data):
    set_background(slide)
    add_left_bar(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8.5), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = data.get("title", "")
    p.font.name = THEME.FONT_JP
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME.PRIMARY
    subtitle = data.get("subtitle", "")
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(0.4))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = THEME.FONT_EN
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = THEME.SECONDARY

def create_list(slide, data):
    add_header(slide, data.get("title", ""), data.get("subtitle", ""))
    items = data.get("items", [])
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = THEME.FONT_JP
        p.font.size = Pt(18)
        p.font.color.rgb = THEME.TEXT
        p.space_after = Pt(12)

def create_three_points_circle(slide, data, theme=THEME):
    add_header(slide, data.get('title', ''), data.get('subtitle', ''))
    items = data.get('items', [])
    start_y = Inches(1.5)
    col_width = Inches(2.8)
    gap = Inches(0.4)
    start_x = Inches(0.7)
    for i, item in enumerate(items[:3]):
        x = start_x + i * (col_width + gap)
        circle_x = x + Inches(0.8)
        circle_y = start_y
        circle_size = Inches(1.4)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, circle_size, circle_size)
        circle.fill.background()
        circle.line.color.rgb = theme.PRIMARY
        circle.line.width = Pt(2)
        text_height = Inches(0.55)
        text_y = circle_y + (circle_size - text_height) / 2 + Inches(0.02)
        num_box = slide.shapes.add_textbox(circle_x, text_y, circle_size, text_height)
        tf = num_box.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        head_box = slide.shapes.add_textbox(x, start_y + Inches(1.55), col_width, Inches(0.4))
        head_box.text_frame.word_wrap = True
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        p.alignment = PP_ALIGN.CENTER
        if desc:
            desc_box = slide.shapes.add_textbox(x, start_y + Inches(1.95), col_width, Inches(1.5))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(10)
            p.font.color.rgb = theme.TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER

def create_point_list_compact(slide, data, theme=THEME):
    add_header(slide, data.get('title', ''), data.get('subtitle', ''))
    items = data.get('items', [])
    start_y = Inches(1.4)
    item_height = Inches(0.85)
    for i, item in enumerate(items[:4]):
        y_pos = start_y + i * item_height
        num_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(0.5), Inches(0.5))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{i+1:02d}"
        p.font.name = theme.FONT_EN
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = theme.PRIMARY
        
        lines = item.split('\n')
        heading = lines[0] if lines else ''
        desc = '\n'.join(lines[1:]) if len(lines) > 1 else ''
        
        head_box = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(8.0), Inches(0.35))
        p = head_box.text_frame.paragraphs[0]
        p.text = heading
        p.font.name = theme.FONT_JP
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = theme.TEXT
        
        if desc:
            desc_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.35), Inches(8.0), Inches(0.45))
            desc_box.text_frame.word_wrap = True
            p = desc_box.text_frame.paragraphs[0]
            p.text = desc
            p.font.name = theme.FONT_JP
            p.font.size = Pt(11)
            p.font.color.rgb = theme.TEXT_LIGHT

TEMPLATES = {
    "cover": create_cover,
    "list": create_list,
    "three_points_circle": create_three_points_circle,
    "point_list_compact": create_point_list_compact,
}

SLIDES_TO_USE = [1, 2, 9, 10]

SLIDE_CONTENT = {
    1: {
        "template": "cover",
        "title": "量子コンピュータの基礎と最新動向",
        "subtitle": "Introduction to Quantum Computing",
    },
    2: {
        "template": "list",
        "title": "本日のアジェンダ",
        "subtitle": "Agenda",
        "items": [
            "量子計算を支える3つの基盤原理",
            "従来のコンピュータ（古典情報）との決定的な違い",
            "産業界に変革をもたらす主な応用分野",
            "実用化（社会実装）に向けた今後の課題とロードマップ",
        ],
    },
    9: {
        "template": "three_points_circle",
        "title": "量子計算を可能にする3つの基盤原理",
        "subtitle": "Core Principles",
        "items": [
            "量子ビット\n情報の最小単位。0と1だけでなく、その中間の状態も同時に保持できる。",
            "重ね合わせ\n複数の状態が同時に存在する性質。膨大な計算の並列処理を可能にする。",
            "量子もつれ\n離れた粒子同士が互いに影響し合う相関関係。高速な情報伝達と制御に不可欠。"
        ]
    },
    10: {
        "template": "point_list_compact",
        "title": "実用化へのロードマップ",
        "subtitle": "Roadmap",
        "items": [
            "理論構築・基礎研究\n量子アルゴリズムおよびハードウェアの基本理論の確立",
            "小規模な実証実験\nノイズを含む中規模な量子デバイス（NISQ）による検証",
            "特定分野での商用利用\n化学計算や最適化問題など、特定領域での先行実用化",
            "汎用量子計算の実現\n量子誤り訂正技術の確立による、完全な汎用マシンの普及"
        ]
    },
}

def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    for slide_no in SLIDES_TO_USE:
        data = SLIDE_CONTENT[slide_no]
        template_name = data["template"]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        TEMPLATES[template_name](slide, data)
    output_path = Path("quantum_presentation.pptx")
    prs.save(output_path)
    print(f"Generation completed: {output_path}")

if __name__ == "__main__":
    generate_presentation()